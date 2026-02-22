from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import io
import os
import uuid

app = Flask(__name__)
CORS(app)

# Fungsi baru untuk membungkus semua section dalam satu daftar XML (p14:sectionLst)
def apply_sections(prs, sections_dict):
    """
    Mengintegrasikan fitur Section PowerPoint 2010+ ke dalam file .pptx
    menggunakan Office 2010/2013+ XML Schema.
    """
    ns_p14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    
    # Pastikan elemen extLst ada di level presentation root
    try:
        ext_lst = prs.element.find(qn('p:extLst'))
        if ext_lst is None:
            ext_lst = prs.element.add_extLst()
    except Exception:
        ext_lst = prs.element.add_extLst()

    # Bangun isi XML untuk tiap-tiap section
    sections_xml = ""
    for name, sld_ids in sections_dict.items():
        section_id = f"{{{str(uuid.uuid4()).upper()}}}"
        # Daftarkan ID slide yang masuk ke dalam section ini
        sld_id_xml = "".join([f'<p14:sldId id="{sid}"/>' for sid in sld_ids])
        
        sections_xml += f'''
            <p14:section name="{name}" id="{section_id}">
                <p14:sldIdLst>
                    {sld_id_xml}
                </p14:sldIdLst>
            </p14:section>'''

    # Bungkus semua section ke dalam kontainer p14:sectionLst
    # URI {521415D9...} adalah standar untuk fitur Section Modern
    full_xml = f'''
    <p:ext xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" 
           uri="{{521415D9-36F0-43E3-9257-96A12269D11F}}">
        <p14:sectionLst xmlns:p14="{ns_p14}">
            {sections_xml}
        </p14:sectionLst>
    </p:ext>
    '''
    
    ext = parse_xml(full_xml)
    ext_lst.append(ext)

@app.route('/')
def home():
    return "<h1>Server Asisten Gereja Aktif!</h1><p>Siap memproses data PowerPoint dengan fitur Modern Section (Office 2021 Compatible).</p>"

@app.route('/generate-ppt', methods=['POST', 'OPTIONS'])
def generate_ppt():
    if request.method == 'OPTIONS':
        return jsonify({"status": "ok"}), 200

    data = request.json
    template_file = 'Source_PowerPoint.pptx'
    
    base_path = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(base_path, template_file)
    
    try:
        prs = Presentation(template_path)
    except Exception as e:
        return jsonify({"error": f"File template tidak ditemukan: {e}"}), 500

    sections_map = {}
    current_section = "Default"

    for item in data.get('slides', []):
        try:
            # Update nama section jika ditemukan field section baru
            if item.get('section'):
                current_section = item['section']
            
            if current_section not in sections_map:
                sections_map[current_section] = []

            layout_idx = int(item.get('layout_idx', 0))
            if layout_idx >= len(prs.slide_layouts):
                layout_idx = 0
                
            layout_dipilih = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(layout_dipilih)
            
            # Catat ID slide untuk didaftarkan ke XML Section nanti
            sections_map[current_section].append(slide.slide_id)
            
            shapes = sorted(slide.placeholders, key=lambda p: p.placeholder_format.idx)

            if len(shapes) > 0 and item.get('judul'):
                shapes[0].text = item['judul']
            
            if len(shapes) > 1 and item.get('isi'):
                tf = shapes[1].text_frame
                tf.text = item['isi']
                
        except Exception as e:
            print(f"Gagal proses slide: {e}")
            continue

    # Terapkan semua section sekaligus ke file PPTX
    if sections_map:
        apply_sections(prs, sections_map)

    target_stream = io.BytesIO()
    prs.save(target_stream)
    target_stream.seek(0)

    return send_file(
        target_stream,
        as_attachment=True,
        download_name="Ibadah_Minggu_Sectioned.pptx",
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 10000))
    app.run(host='0.0.0.0', port=port)
